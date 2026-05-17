import io
import json
import csv as csv_mod
import logging
import re
from typing import Optional
from dataclasses import dataclass

logger = logging.getLogger("HeatAI.crawler")


@dataclass
class ParsedDocument:
    filename: str
    file_ext: str
    source_url: str
    title: str
    text: str
    text_length: int
    parse_status: str  # "success" | "partial" | "failed"


class DocumentParser:

    SUPPORTED_EXTS = {
        ".pdf", ".docx", ".doc", ".html", ".htm", ".shtml", ".shtm", ".txt",
        ".md", ".markdown", ".csv", ".json",
        ".xlsx", ".xls", ".pptx", ".ppt", ".epub",
        ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp",
    }

    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}

    REMOVE_TAGS = [
        "script", "style", "nav", "footer", "header", "aside", "noscript",
        "a", "button", "input", "form", "img", "iframe", "select",
        "textarea", "label", "link", "meta", "figure", "figcaption",
    ]

    _ocr = None

    @classmethod
    def parse(cls, content_bytes: bytes, filename: str, source_url: str = "") -> ParsedDocument:
        ext = cls._get_ext(filename)
        if ext not in cls.SUPPORTED_EXTS:
            raise ValueError(f"不支持的文件类型: {ext}")

        try:
            if ext == ".pdf":
                text, title = cls._parse_pdf(content_bytes, filename)
            elif ext in (".docx", ".doc"):
                text, title = cls._parse_docx(content_bytes, filename)
            elif ext in (".html", ".htm", ".shtml", ".shtm"):
                text, title = cls._parse_html(content_bytes, filename)
            elif ext == ".txt":
                text, title = cls._parse_txt(content_bytes, filename)
            elif ext in (".md", ".markdown"):
                text, title = cls._parse_markdown(content_bytes, filename)
            elif ext == ".csv":
                text, title = cls._parse_csv(content_bytes, filename)
            elif ext == ".json":
                text, title = cls._parse_json(content_bytes, filename)
            elif ext in (".xlsx", ".xls"):
                text, title = cls._parse_xlsx(content_bytes, filename)
            elif ext in (".pptx", ".ppt"):
                text, title = cls._parse_pptx(content_bytes, filename)
            elif ext == ".epub":
                text, title = cls._parse_epub(content_bytes, filename)
            elif ext in cls.IMAGE_EXTS:
                text, title = cls._parse_image(content_bytes, filename)
            else:
                raise ValueError(f"不支持的文件类型: {ext}")

            # 如果 text 为空但 content 有内容, 则为部分成功
            if not text or not text.strip():
                if len(content_bytes) > 100:
                    return ParsedDocument(
                        filename=filename,
                        file_ext=ext,
                        source_url=source_url,
                        title=title or filename,
                        text="(该文档无法提取文本内容，请手动查看)",
                        text_length=0,
                        parse_status="partial",
                    )
                else:
                    return ParsedDocument(
                        filename=filename,
                        file_ext=ext,
                        source_url=source_url,
                        title=title or filename,
                        text="",
                        text_length=0,
                        parse_status="failed",
                    )

            return ParsedDocument(
                filename=filename,
                file_ext=ext,
                source_url=source_url,
                title=title or filename,
                text=text,
                text_length=len(text),
                parse_status="success",
            )
        except Exception as e:
            logger.error(f"[解析失败] {filename}: {e}")
            return ParsedDocument(
                filename=filename,
                file_ext=ext,
                source_url=source_url,
                title=filename,
                text=f"(解析异常: {str(e)})",
                text_length=0,
                parse_status="failed",
            )

    @staticmethod
    def _get_ext(filename: str) -> str:
        if "." in filename:
            return "." + filename.rsplit(".", 1)[-1].lower()
        return ""

    @staticmethod
    def _parse_pdf(content_bytes: bytes, filename: str):
        import pdfplumber

        texts = []
        with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    texts.append(page_text)

        full_text = "\n\n".join(texts)
        total_chars = len(full_text.replace("\n", "").replace(" ", ""))
        total_pages = len(pdf.pages)

        if total_pages > 0 and total_chars / total_pages < 30:
            logger.info(f"[PDF] 字符密度过低 ({total_chars}/{total_pages}页), 尝试OCR")
            ocr = DocumentParser._get_ocr()
            if ocr is not None:
                ocr_texts = []
                for page in pdf.pages:
                    img = page.to_image(resolution=200)
                    ocr_result = ocr.ocr(img.original, cls=True)
                    if ocr_result and ocr_result[0]:
                        page_lines = [line[1][0] for line in ocr_result[0]]
                        ocr_texts.append("\n".join(page_lines))
                if ocr_texts:
                    ocr_full = "\n\n".join(ocr_texts)
                    if len(ocr_full.replace("\n", "").strip()) > total_chars:
                        logger.info(f"[PDF] OCR增强: {total_chars} → {len(ocr_full)} 字符")
                        full_text = ocr_full

        return full_text, DocumentParser._guess_title(full_text, filename)

    @staticmethod
    def _parse_docx(content_bytes: bytes, filename: str):
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(content_bytes))
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        for table in doc.tables:
            table_lines = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_lines.append(" | ".join(cells))
            if table_lines:
                paragraphs.append("\n" + "\n".join(table_lines) + "\n")

        full_text = "\n\n".join(paragraphs)
        return full_text, DocumentParser._guess_title(full_text, filename)

    @staticmethod
    def _parse_html(content_bytes: bytes, filename: str):
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content_bytes, "lxml")

        for tag in soup(DocumentParser.REMOVE_TAGS):
            tag.decompose()

        title_tag = soup.find("title")
        title = title_tag.get_text(strip=True) if title_tag else filename.rsplit(".", 1)[0]

        body = soup.find("body")
        if body:
            text = body.get_text(separator="\n", strip=True)
        else:
            text = soup.get_text(separator="\n", strip=True)

        return text, title

    @staticmethod
    def _parse_txt(content_bytes: bytes, filename: str):
        text = content_bytes.decode("utf-8", errors="replace")
        return text, filename.rsplit(".", 1)[0]

    @staticmethod
    def _guess_title(text: str, filename: str) -> str:
        if not text:
            return filename
        lines = text.strip().split("\n")
        for line in lines[:10]:
            line = line.strip()
            if line and len(line) >= 5 and len(line) <= 100:
                return line
        return filename

    @staticmethod
    def _parse_markdown(content_bytes: bytes, filename: str):
        text = content_bytes.decode("utf-8", errors="replace")
        text = re.sub(r"```[\s\S]*?```", lambda m: f"\n[代码块]\n{m.group(0)}\n[/代码块]\n", text)
        text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"[图片: \1]", text)
        text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
        text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
        text = re.sub(r"\*{1,3}([^*]+)\*{1,3}", r"\1", text)
        text = re.sub(r"`([^`]+)`", r"\1", text)
        text = re.sub(r"^\s*[-*+]\s+", "• ", text, flags=re.MULTILINE)
        text = re.sub(r"^\s*\d+\.\s+", "", text, flags=re.MULTILINE)
        return text, DocumentParser._guess_title(text, filename)

    @staticmethod
    def _parse_csv(content_bytes: bytes, filename: str):
        text = content_bytes.decode("utf-8", errors="replace")
        reader = csv_mod.reader(io.StringIO(text))
        rows = list(reader)
        if not rows:
            return "", filename

        max_cols = max(len(row) for row in rows)
        col_widths = [0] * max_cols
        for row in rows:
            for i, c in enumerate(row):
                col_widths[i] = max(col_widths[i], len(str(c)))

        lines = []
        for ri, row in enumerate(rows):
            padded = [str(c).ljust(col_widths[i]) for i, c in enumerate(row)]
            lines.append(" | ".join(padded))
            if ri == 0:
                lines.append(" | ".join("-" * col_widths[i] for i in range(max_cols)))

        return "\n".join(lines), DocumentParser._guess_title(rows[0][0] if rows[0] else "", filename)

    @staticmethod
    def _parse_json(content_bytes: bytes, filename: str):
        text = content_bytes.decode("utf-8", errors="replace")
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            return text, filename

        lines = []

        def _flatten(obj, prefix: str = "", depth: int = 0):
            if depth > 10:
                lines.append(f"{prefix}{str(obj)[:200]}")
                return
            if isinstance(obj, dict):
                for k, v in obj.items():
                    key = f"{prefix}.{k}" if prefix else k
                    if isinstance(v, (dict, list)):
                        lines.append(f"## {key}")
                        _flatten(v, key, depth + 1)
                    else:
                        lines.append(f"- {key}: {v}")
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    label = f"{prefix}[{i}]"
                    if isinstance(item, (dict, list)):
                        lines.append(f"## {label}")
                        _flatten(item, label, depth + 1)
                    else:
                        lines.append(f"- {label}: {item}")

        _flatten(data)
        return "\n".join(lines), DocumentParser._guess_title("\n".join(lines), filename)

    @staticmethod
    def _parse_xlsx(content_bytes: bytes, filename: str):
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
        all_texts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_texts.append(f"## Sheet: {sheet_name}")

            rows_data = []
            for row in ws.iter_rows(values_only=True):
                if row and any(c is not None for c in row):
                    rows_data.append([str(c) if c is not None else "" for c in row])

            if not rows_data:
                all_texts.append("(空工作表)")
                continue

            max_cols = max(len(row) for row in rows_data)
            col_widths = [0] * max_cols
            for row in rows_data[:50]:
                for i, c in enumerate(row):
                    col_widths[i] = max(col_widths[i], min(len(c), 40))

            for ri, row in enumerate(rows_data[:200]):
                padded = [str(c)[:40].ljust(col_widths[i]) for i, c in enumerate(row)]
                all_texts.append(" | ".join(padded))
                if ri == 0:
                    all_texts.append(" | ".join("-" * min(col_widths[i], 40) for i in range(max_cols)))

            if len(rows_data) > 200:
                all_texts.append(f"... (省略 {len(rows_data) - 200} 行)")
            all_texts.append("")

        wb.close()
        return "\n".join(all_texts), DocumentParser._guess_title("\n".join(all_texts), filename)

    @staticmethod
    def _parse_pptx(content_bytes: bytes, filename: str):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content_bytes))
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_parts = [f"## 幻灯片 {i + 1}"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)

                if shape.has_table:
                    table = shape.table
                    table_lines = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        table_lines.append(" | ".join(cells))
                    if table_lines:
                        slide_parts.append("\n" + "\n".join(table_lines) + "\n")

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"\n[备注]\n{notes}")

            slides_text.append("\n".join(slide_parts))

        return "\n\n".join(slides_text), DocumentParser._guess_title("\n\n".join(slides_text), filename)

    @staticmethod
    def _parse_epub(content_bytes: bytes, filename: str):
        from ebooklib import epub
        from bs4 import BeautifulSoup

        book = epub.read_epub(io.BytesIO(content_bytes))
        chapters = []

        titles = book.get_metadata("DC", "title")
        epub_title = titles[0][0] if titles else filename

        for item in book.get_items_of_type(9):
            content = item.get_content().decode("utf-8", errors="replace")
            soup = BeautifulSoup(content, "lxml")

            for tag in soup(DocumentParser.REMOVE_TAGS):
                tag.decompose()

            text = soup.get_text(separator="\n", strip=True)
            if text and len(text.strip()) > 20:
                chapters.append(text)

        if not chapters:
            return "", epub_title

        return "\n\n".join(chapters), epub_title

    @classmethod
    def _get_ocr(cls):
        if cls._ocr is not None:
            return cls._ocr
        try:
            from paddleocr import PaddleOCR
            cls._ocr = PaddleOCR(lang="ch", use_angle_cls=True, show_log=False)
            logger.info("[OCR] PaddleOCR 引擎初始化成功")
            return cls._ocr
        except ImportError:
            logger.warning("[OCR] PaddleOCR 未安装, 图片/扫描件PDF将被跳过")
            return None
        except Exception as e:
            logger.warning(f"[OCR] PaddleOCR 初始化失败: {e}")
            return None

    @staticmethod
    def _parse_image(content_bytes: bytes, filename: str):
        ocr = DocumentParser._get_ocr()
        if ocr is None:
            return "(OCR引擎未就绪)", filename

        from PIL import Image
        img = Image.open(io.BytesIO(content_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        import numpy as np
        img_array = np.array(img)

        result = ocr.ocr(img_array, cls=True)
        if not result or not result[0]:
            return "", filename

        lines = [line[1][0] for line in result[0]]
        text = "\n".join(lines)
        return text, DocumentParser._guess_title(text, filename)


parser = DocumentParser()
