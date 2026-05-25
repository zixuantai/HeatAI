import io
import json
import csv as csv_mod
import logging
import re
from typing import List, Tuple

logger = logging.getLogger(__name__)


class DocumentParser:
    SUPPORTED_TYPES = {
        "pdf", "docx", "doc", "html", "htm", "txt",
        "md", "markdown", "csv", "json",
        "xlsx", "xls", "pptx", "ppt", "epub",
        "png", "jpg", "jpeg", "bmp", "tiff", "webp",
    }

    IMAGE_TYPES = {"png", "jpg", "jpeg", "bmp", "tiff", "webp"}

    REMOVE_TAGS = [
        "script", "style", "nav", "footer", "header", "aside", "noscript",
        "a", "button", "input", "form", "img", "iframe", "select",
        "textarea", "label", "link", "meta", "figure", "figcaption",
    ]

    _ocr = None

    @staticmethod
    def parse(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext not in DocumentParser.SUPPORTED_TYPES:
            raise ValueError(f"不支持的文件类型: .{ext}")

        try:
            if ext == "pdf":
                return DocumentParser._parse_pdf(file_bytes, filename)
            elif ext in ("docx", "doc"):
                return DocumentParser._parse_doc(file_bytes, filename)
            elif ext in ("html", "htm"):
                return DocumentParser._parse_html(file_bytes, filename)
            elif ext == "txt":
                return DocumentParser._parse_txt(file_bytes, filename)
            elif ext in ("md", "markdown"):
                return DocumentParser._parse_markdown(file_bytes, filename)
            elif ext == "csv":
                return DocumentParser._parse_csv(file_bytes, filename)
            elif ext == "json":
                return DocumentParser._parse_json(file_bytes, filename)
            elif ext in ("xlsx", "xls"):
                return DocumentParser._parse_xlsx(file_bytes, filename)
            elif ext in ("pptx", "ppt"):
                return DocumentParser._parse_pptx(file_bytes, filename)
            elif ext == "epub":
                return DocumentParser._parse_epub(file_bytes, filename)
            elif ext in DocumentParser.IMAGE_TYPES:
                return DocumentParser._parse_image(file_bytes, filename)
            else:
                raise ValueError(f"不支持的文件类型: .{ext}")
        except ValueError:
            raise
        except MemoryError:
            raise ValueError("文件解析时内存不足，文件可能过大")
        except Exception as e:
            ext_hint = f"文件扩展名为 .{ext}，" if ext else ""
            raise ValueError(f"{ext_hint}文件解析失败: {type(e).__name__}: {e}")

    @staticmethod
    def _parse_pdf(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        import pdfplumber

        try:
            pdf = pdfplumber.open(io.BytesIO(file_bytes))
        except Exception as e:
            raise ValueError(f"PDF文件无法打开，可能已损坏或不是有效的PDF格式: {e}")

        try:
            texts: List[str] = []
            for page in pdf.pages:
                page_parts: List[str] = []

                page_text = page.extract_text()
                if page_text:
                    page_parts.append(page_text)

                tables = page.extract_tables()
                if tables:
                    page_parts.append("")
                    for ti, table in enumerate(tables):
                        if not table:
                            continue
                        page_parts.append(f"[表格 {ti + 1}]")

                        filtered_rows = []
                        for row in table:
                            filtered = [str(c).strip() if c is not None else "" for c in row]
                            if any(filtered):
                                filtered_rows.append(filtered)
                        if not filtered_rows:
                            continue

                        header_row = filtered_rows[0]
                        col_count = len(header_row)
                        header_str = " | ".join(header_row)
                        sep_str = " | ".join("---" for _ in range(col_count))
                        page_parts.append(header_str)
                        page_parts.append(sep_str)

                        for row in filtered_rows[1:]:
                            padded = row + [""] * (col_count - len(row))
                            page_parts.append(" | ".join(padded[:col_count]))
                        page_parts.append("")

                if page_parts:
                    texts.append("\n".join(page_parts))

            full_text = "\n\n---\n\n".join(texts)
            total_chars = len(full_text.replace("\n", "").replace(" ", "").replace("-", ""))

            total_pages = len(pdf.pages)
            if total_pages > 0 and total_chars / total_pages < 30:
                logger.info(f"[PDF] 字符密度过低 ({total_chars}/{total_pages}页), 尝试OCR")
                ocr = DocumentParser._get_ocr()
                if ocr is not None:
                    ocr_texts: List[str] = []
                    for page in pdf.pages:
                        img = page.to_image(resolution=200)
                        ocr_result = ocr.ocr(img.original, cls=True)
                        if ocr_result and ocr_result[0]:
                            page_lines = [line[1][0] for line in ocr_result[0]]
                            ocr_texts.append("\n".join(page_lines))
                    if ocr_texts:
                        ocr_full = "\n\n---\n\n".join(ocr_texts)
                        if len(ocr_full.replace("\n", "").strip()) > total_chars:
                            logger.info(f"[PDF] OCR增强完成: {total_chars} → {len(ocr_full)} 字符")
                            full_text = ocr_full

            title = filename.rsplit(".", 1)[0]
            if not full_text.strip():
                raise ValueError("PDF文件未能提取到任何文字内容，可能为扫描件且OCR未启用")
            return full_text, title
        finally:
            try:
                pdf.close()
            except Exception:
                pass

    @staticmethod
    def _parse_doc(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext == "docx":
            return DocumentParser._parse_docx(file_bytes, filename)
        else:
            return DocumentParser._parse_old_doc(file_bytes, filename)

    @staticmethod
    def _parse_docx(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(file_bytes))
        paragraphs: List[str] = []

        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            if "Heading" in style or "heading" in style or "标题" in style:
                paragraphs.append(f"## {text}")
            else:
                paragraphs.append(text)

        for table in doc.tables:
            table_lines: List[str] = []
            rows_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if not rows_data:
                continue
            header = rows_data[0]
            col_count = len(header)
            table_lines.append(" | ".join(header))
            table_lines.append(" | ".join("---" for _ in range(col_count)))
            for row in rows_data[1:]:
                padded = row + [""] * (col_count - len(row))
                table_lines.append(" | ".join(padded[:col_count]))
            if table_lines:
                paragraphs.append("\n" + "\n".join(table_lines) + "\n")

        full_text = "\n\n".join(paragraphs)
        title = filename.rsplit(".", 1)[0]
        return full_text, title

    @staticmethod
    def _parse_old_doc(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        title = filename.rsplit(".", 1)[0]

        sig = file_bytes[:8]
        if sig[:2] != b"\xd0\xcf":
            try:
                return DocumentParser._parse_docx(file_bytes, filename)
            except Exception:
                raise ValueError("无法识别该 .doc 文件格式，请将其另存为 .docx 格式后再上传")

        try:
            import olefile
        except ImportError:
            raise ValueError(
                "解析旧版 .doc 格式需要 olefile 库，请运行: pip install olefile，"
                "或将该文件另存为 .docx 格式后再上传"
            )

        try:
            ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        except Exception:
            raise ValueError("无法打开 .doc 文件，文件可能已损坏，请将其另存为 .docx 格式后再上传")

        try:
            word_stream = ole.openstream("WordDocument")
            word_bytes = word_stream.read()
        except Exception:
            ole.close()
            raise ValueError("无法读取 .doc 文件内容，请将其另存为 .docx 格式后再上传")

        text_parts: List[str] = []

        try:
            fc_min = int.from_bytes(word_bytes[0x18:0x1C], "little")
            piece_table_offset = int.from_bytes(word_bytes[0x1A2:0x1A6], "little")
            base_offset = fc_min

            piece_table_data = word_bytes[piece_table_offset:]
            piece_count = (piece_table_data[0] | (piece_table_data[1] << 8)) + 1

            pos = 2 + piece_count * 4 + 2
            cp_start = 0
            for i in range(piece_count - 1):
                cp_end = int.from_bytes(piece_table_data[2 + i * 4:6 + i * 4], "little")

                desc = int.from_bytes(piece_table_data[pos:pos + 8], "little")
                pos += 8

                fc_value = desc & 0x3FFFFFFF
                is_unicode = not bool(desc & 0x40000000)

                length = cp_end - cp_start
                offset = (fc_value - base_offset) if fc_value >= base_offset else fc_value

                if offset >= 0 and offset + length * (2 if is_unicode else 1) <= len(word_bytes):
                    if is_unicode:
                        chunk = word_bytes[offset:offset + length * 2].decode("utf-16-le", errors="replace")
                    else:
                        enc = "gbk"
                        try:
                            word_bytes[offset:offset + length].decode("cp1252")
                            enc = "cp1252"
                        except Exception:
                            pass
                        chunk = word_bytes[offset:offset + length].decode(enc, errors="replace")
                    text_parts.append(chunk)

                cp_start = cp_end

            if not text_parts:
                full_text = word_bytes.decode("utf-16-le", errors="replace")
                printable = "".join(c for c in full_text if c.isprintable() or c in "\n\r\t")
                if len(printable) < 10:
                    raise ValueError("无法从 .doc 文件中提取有用文本，请另存为 .docx 格式后再上传")
                return printable, title

            full_text = "".join(text_parts)
            if not full_text.strip():
                raise ValueError("该 .doc 文件未能提取到文字内容，请另存为 .docx 格式后再上传")

            return full_text, title
        finally:
            ole.close()

    @staticmethod
    def _parse_html(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        from bs4 import BeautifulSoup

        encoding = DocumentParser._detect_encoding(file_bytes)
        logger.info(f"[HTML解析] 检测到编码: {encoding}, 文件名: {filename}")

        html_content = file_bytes.decode(encoding, errors="replace")

        meta_encoding = DocumentParser._extract_html_charset(html_content)
        if meta_encoding and meta_encoding.lower() != encoding.lower():
            try:
                html_content = file_bytes.decode(meta_encoding, errors="replace")
                logger.info(f"[HTML解析] 使用HTML声明的编码: {meta_encoding}")
            except (UnicodeDecodeError, LookupError):
                logger.warning(f"[HTML解析] HTML声明编码 {meta_encoding} 无效, 回退到 {encoding}")

        soup = BeautifulSoup(html_content, "lxml")

        for tag in soup(DocumentParser.REMOVE_TAGS):
            tag.decompose()

        title_tag = soup.find("title")
        title = title_tag.get_text(strip=True) if title_tag else filename.rsplit(".", 1)[0]

        body = soup.find("body")
        if body:
            text = body.get_text(separator="\n", strip=True)
        else:
            text = soup.get_text(separator="\n", strip=True)

        if not text or not text.strip():
            text = soup.get_text(separator="\n")
            text = "\n".join(line.strip() for line in text.split("\n") if line.strip())

        if not text or not text.strip():
            logger.warning(f"[HTML解析] 文档解析后内容为空，可能为SPA动态渲染页面: {filename}")
            return "", title

        return text, title

    @staticmethod
    def _extract_html_charset(html: str) -> str | None:
        import re
        meta_pattern = re.compile(
            r'<meta[^>]+charset\s*=\s*["\']?([\w\-\d]+)["\']?',
            re.IGNORECASE,
        )
        match = meta_pattern.search(html[:4096])
        if match:
            return match.group(1)
        return None

    @staticmethod
    def _detect_encoding(file_bytes: bytes) -> str:
        try:
            import chardet
            result = chardet.detect(file_bytes)
            detected = result.get("encoding")
            confidence = result.get("confidence", 0)
            if detected and confidence > 0.7:
                try:
                    file_bytes.decode(detected)
                    return detected
                except (UnicodeDecodeError, LookupError):
                    pass
        except ImportError:
            pass

        for enc in ["utf-8", "gb18030", "gbk", "gb2312", "latin-1"]:
            try:
                file_bytes.decode(enc)
                return enc
            except (UnicodeDecodeError, LookupError):
                continue
        return "utf-8"

    @staticmethod
    def _parse_txt(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        encoding = DocumentParser._detect_encoding(file_bytes)
        text = file_bytes.decode(encoding, errors="replace")
        title = filename.rsplit(".", 1)[0]
        return text, title

    @staticmethod
    def _parse_markdown(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        encoding = DocumentParser._detect_encoding(file_bytes)
        md_text = file_bytes.decode(encoding, errors="replace")
        title = filename.rsplit(".", 1)[0]

        md_text = re.sub(r"```[\s\S]*?```", lambda m: f"\n[代码块]\n{m.group(0)}\n[/代码块]\n", md_text)
        md_text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"[图片: \1]", md_text)
        md_text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", md_text)
        md_text = re.sub(r"^#{1,6}\s+", "", md_text, flags=re.MULTILINE)
        md_text = re.sub(r"\*{1,3}([^*]+)\*{1,3}", r"\1", md_text)
        md_text = re.sub(r"`([^`]+)`", r"\1", md_text)
        md_text = re.sub(r"^\s*[-*+]\s+", "• ", md_text, flags=re.MULTILINE)
        md_text = re.sub(r"^\s*\d+\.\s+", "", md_text, flags=re.MULTILINE)
        md_text = re.sub(r"\|", " ", md_text)
        md_text = re.sub(r"^\s*[-:]+\s*$", "", md_text, flags=re.MULTILINE)

        return md_text, title

    @staticmethod
    def _parse_csv(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        encoding = DocumentParser._detect_encoding(file_bytes)
        text = file_bytes.decode(encoding, errors="replace")
        title = filename.rsplit(".", 1)[0]

        lines: List[str] = []
        reader = csv_mod.reader(io.StringIO(text))
        rows = list(reader)
        if not rows:
            return "", title

        max_cols = max(len(row) for row in rows)
        col_widths: List[int] = [0] * max_cols
        for row in rows:
            for i, cell in enumerate(row):
                col_widths[i] = max(col_widths[i], len(str(cell)))

        for ri, row in enumerate(rows):
            padded = [str(c).ljust(col_widths[i]) for i, c in enumerate(row)]
            lines.append(" | ".join(padded))
            if ri == 0:
                lines.append(" | ".join("-" * col_widths[i] for i in range(max_cols)))

        return "\n".join(lines), title

    @staticmethod
    def _parse_json(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        encoding = DocumentParser._detect_encoding(file_bytes)
        text = file_bytes.decode(encoding, errors="replace")
        title = filename.rsplit(".", 1)[0]

        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            return text, title

        lines: List[str] = []

        def _flatten(obj, prefix: str = "", depth: int = 0) -> None:
            if depth > 10:
                lines.append(f"{prefix}{str(obj)[:200]}")
                return
            if isinstance(obj, dict):
                for k, v in obj.items():
                    key = f"{prefix}.{k}" if prefix else k
                    if isinstance(v, (dict, list)):
                        lines.append(f"## {key}")
                        _flatten(v, key, depth + 1)
                    else:
                        lines.append(f"- {key}: {v}")
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    label = f"{prefix}[{i}]"
                    if isinstance(item, (dict, list)):
                        lines.append(f"## {label}")
                        _flatten(item, label, depth + 1)
                    else:
                        lines.append(f"- {label}: {item}")

        _flatten(data)
        return "\n".join(lines), title

    @staticmethod
    def _parse_xlsx(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        from openpyxl import load_workbook

        title = filename.rsplit(".", 1)[0]
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        all_texts: List[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_texts.append(f"## Sheet: {sheet_name}")

            rows_data: List[List[str]] = []
            for row in ws.iter_rows(values_only=True):
                if row and any(c is not None for c in row):
                    rows_data.append([str(c) if c is not None else "" for c in row])

            if not rows_data:
                all_texts.append("(空工作表)")
                continue

            max_cols = max(len(row) for row in rows_data)
            col_widths = [0] * max_cols
            for row in rows_data[:50]:
                for i, c in enumerate(row):
                    col_widths[i] = max(col_widths[i], min(len(c), 40))

            for ri, row in enumerate(rows_data[:200]):
                padded = [str(c)[:40].ljust(col_widths[i]) for i, c in enumerate(row)]
                all_texts.append(" | ".join(padded))
                if ri == 0:
                    all_texts.append(" | ".join("-" * min(col_widths[i], 40) for i in range(max_cols)))

            if len(rows_data) > 200:
                all_texts.append(f"... (省略 {len(rows_data) - 200} 行)")

            all_texts.append("")

        wb.close()
        return "\n".join(all_texts), title

    @staticmethod
    def _parse_pptx(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        from pptx import Presentation

        title = filename.rsplit(".", 1)[0]
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text: List[str] = []

        for i, slide in enumerate(prs.slides):
            slide_parts: List[str] = [f"## 幻灯片 {i + 1}"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)

                if shape.has_table:
                    table = shape.table
                    table_lines: List[str] = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        table_lines.append(" | ".join(cells))
                    if table_lines:
                        slide_parts.append("\n" + "\n".join(table_lines) + "\n")

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"\n[备注]\n{notes}")

            slides_text.append("\n".join(slide_parts))

        return "\n\n".join(slides_text), title

    @staticmethod
    def _parse_epub(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        from ebooklib import epub
        from bs4 import BeautifulSoup

        title = filename.rsplit(".", 1)[0]
        book = epub.read_epub(io.BytesIO(file_bytes))
        chapters: List[str] = []

        epub_title = title
        titles = book.get_metadata("DC", "title")
        if titles:
            epub_title = titles[0][0]

        for item in book.get_items_of_type(9):
            content = item.get_content().decode("utf-8", errors="replace")
            soup = BeautifulSoup(content, "lxml")

            for tag in soup(DocumentParser.REMOVE_TAGS):
                tag.decompose()

            text = soup.get_text(separator="\n", strip=True)
            if text and len(text.strip()) > 20:
                chapters.append(text)

        if not chapters:
            return "", epub_title

        return "\n\n".join(chapters), epub_title

    @classmethod
    def _get_ocr(cls):
        if cls._ocr is not None:
            return cls._ocr
        try:
            from paddleocr import PaddleOCR
            cls._ocr = PaddleOCR(lang="ch", use_angle_cls=True, show_log=False)
            logger.info("[OCR] PaddleOCR 引擎初始化成功")
            return cls._ocr
        except ImportError:
            logger.warning("[OCR] PaddleOCR 未安装, 图片/扫描件PDF将被跳过. 安装: pip install paddleocr")
            return None
        except Exception as e:
            logger.warning(f"[OCR] PaddleOCR 初始化失败: {e}")
            return None

    @staticmethod
    def _parse_image(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        ocr = DocumentParser._get_ocr()
        if ocr is None:
            raise ValueError(
                "OCR 引擎未就绪，无法解析图片文件。请安装: pip install paddleocr"
            )

        from PIL import Image

        img = Image.open(io.BytesIO(file_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        import numpy as np
        img_array = np.array(img)

        result = ocr.ocr(img_array, cls=True)
        if not result or not result[0]:
            return "", filename.rsplit(".", 1)[0]

        lines = [line[1][0] for line in result[0]]
        text = "\n".join(lines)
        title = filename.rsplit(".", 1)[0]
        return text, title


document_parser = DocumentParser()
